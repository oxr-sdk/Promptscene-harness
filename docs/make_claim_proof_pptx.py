# -*- coding: utf-8 -*-
"""
Create a clean PPTX from OXR-SDK-AI-ready-claim-proof-slides.md.

Palette: deep blue, white, cream only.
Fonts: title = Pretendard Bold style, body = Pretendard Light style.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


BLUE = RGBColor(0x0B, 0x2A, 0x6F)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CREAM = RGBColor(0xF2, 0xF0, 0xE8)

TITLE_FONT = "Pretendard"
BODY_FONT = "Pretendard Light"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def add_slide(bg=WHITE):
    slide = prs.slides.add_slide(BLANK)
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = bg
    return slide


def set_font(run, size, color=BLUE, bold=False, font=BODY_FONT):
    run.font.name = TITLE_FONT if bold else font
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold


def text(slide, value, x, y, w, h, size=18, color=BLUE, bold=False,
         align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, font=BODY_FONT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = value
    set_font(run, size, color, bold, font)
    return box


def bullets(slide, items, x, y, w, h, size=18, color=BLUE, bullet="•"):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.space_after = Pt(7)
        p.line_spacing = 1.12
        run = p.add_run()
        run.text = f"{bullet} {item}"
        set_font(run, size, color, False)
    return box


def rect(slide, x, y, w, h, fill=CREAM, line=BLUE, radius=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1)
    return shape


def line(slide, x, y, w, h=0.02, fill=BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape


def title_slide(title, subtitle):
    slide = add_slide(BLUE)
    text(slide, title, 0.9, 2.05, 11.5, 1.45, 54, WHITE, True)
    line(slide, 0.9, 3.75, 3.1, 0.045, WHITE)
    text(slide, subtitle, 0.9, 4.05, 10.6, 0.55, 20, WHITE)
    text(slide, "OXR Research · KISTI · 2026-06-16", 0.9, 6.75, 11.5, 0.3, 11, WHITE)


def header(slide, n, title):
    line(slide, 0, 0, 13.333, 0.08, BLUE)
    text(slide, f"{n:02d}", 0.65, 0.55, 0.7, 0.38, 16, BLUE, True)
    text(slide, title, 1.45, 0.46, 10.8, 0.55, 28, BLUE, True)
    line(slide, 0.65, 1.22, 12.05, 0.018, CREAM)


def big_quote(slide, quote, note=None):
    rect(slide, 0.95, 2.0, 11.45, 2.25, BLUE, BLUE)
    text(slide, quote, 1.25, 2.32, 10.85, 1.08, 28, WHITE, True, valign=MSO_ANCHOR.MIDDLE)
    if note:
        text(slide, note, 1.25, 3.55, 10.85, 0.38, 15, WHITE)


def two_columns(slide, left_title, left_items, right_title, right_items):
    rect(slide, 0.75, 1.65, 5.85, 5.25, CREAM, CREAM)
    rect(slide, 6.75, 1.65, 5.85, 5.25, BLUE, BLUE)
    text(slide, left_title, 1.05, 1.95, 5.25, 0.42, 22, BLUE, True)
    bullets(slide, left_items, 1.05, 2.65, 5.15, 3.75, 17, BLUE)
    text(slide, right_title, 7.05, 1.95, 5.25, 0.42, 22, WHITE, True)
    bullets(slide, right_items, 7.05, 2.65, 5.15, 3.75, 17, WHITE)


def table(slide, rows, x, y, w, h, widths, header_fill=BLUE, font_size=13):
    n_rows = len(rows)
    row_h = h / n_rows
    col_x = [x]
    for cw in widths[:-1]:
        col_x.append(col_x[-1] + w * cw)
    for r, row in enumerate(rows):
        fill = header_fill if r == 0 else (CREAM if r % 2 else WHITE)
        txt_color = WHITE if r == 0 else BLUE
        yy = y + r * row_h
        rect(slide, x, yy, w, row_h, fill, fill)
        cx = x
        for c, cell in enumerate(row):
            cw = w * widths[c]
            text(slide, cell, cx + 0.12, yy + 0.11, cw - 0.2, row_h - 0.16,
                 font_size if r else font_size + 1, txt_color, r == 0)
            cx += cw


# 1 Cover
title_slide("OXR-SDK는 AI-Ready 되었는가", "정적 샘플에서 재생성 가능한 프레임워크로")

# 2 Claim
sl = add_slide(WHITE)
header(sl, 1, "주장")
big_quote(sl, "OXR-SDK의 AI-ready 가능성을 XRCollabDemo 위에서 검증했다")
bullets(sl, [
    "검증 무대: 로비 · 룸 생성 · 네트워크 스폰 · 아바타 · 씬 전환이 모두 포함된 XRCollabDemo",
    "핵심 질문: OXR-SDK 위에서 XRCollabDemo 수준의 협업 룸을 정적 샘플 클론 없이 재생성할 수 있는가?",
    "결론: 문서 · 프롬프트 · 검증 기준으로 실제 작동함을 확인",
], 1.05, 4.75, 11.4, 1.7, 17)

# 3 Definition
sl = add_slide(CREAM)
header(sl, 2, "AI-Ready 정의")
big_quote(sl, "AI가 사람의 암묵지 없이 작동 산출물을 만들고 검증할 수 있는 상태")
bullets(sl, [
    "AI가 읽고 실행할 수 있는 절차",
    "반드시 지켜야 하는 작동 조건",
    "결과를 판정할 수 있는 검증 기준",
    "실행 샘플이 아니라 재생성 가능한 레퍼런스",
], 1.25, 4.78, 11.0, 1.75, 18)

# 4 Before After
sl = add_slide(WHITE)
header(sl, 3, "AI-Ready 전과 후")
two_columns(
    sl,
    "기존 샘플 방식",
    ["XRCollabDemo 프로젝트 다운로드", "지식은 씬/프리팹 안에 암묵적으로 포함", "AI는 기존 프로젝트를 열어 보고 따라 함", "실패 원인은 사람이 디버깅"],
    "AI-ready 방식",
    ["OXR-SDK + 문서 + 프롬프트 + 검증 기준", "절차 · 불변식 · 씬 토폴로지로 외부화", "AI가 새 룸을 조립 · 빌드 · 검증", "조건 위반 여부로 판정"],
)

# 5 Components
sl = add_slide(WHITE)
header(sl, 4, "OXR-SDK 구성요소")
table(sl, [
    ["구성요소", "역할"],
    ["com.oxr.sdk", "인증, 에셋 스토어, 파일 스토리지, 룸 관리"],
    ["XumLobby", "로비 UI, 룸 생성, 룸 진입, 룸 씬 전환"],
    ["XumNet", "네트워크 오브젝트 생성과 스폰"],
    ["UnifiedXRMotion", "플랫폼별 아바타와 모션"],
], 0.9, 1.7, 11.55, 3.65, [0.28, 0.72], font_size=13)
text(sl, "검증 관점: 데모 하나가 실행되는가가 아니라, 이 레이어들이 AI가 읽고 조립할 수 있는 형태인가", 0.9, 6.0, 11.6, 0.45, 20, BLUE, True)

# 6 Method
sl = add_slide(CREAM)
header(sl, 5, "검증 방법")
cols = [
    ("1", "문서 확인", "OXR-SDK 문서와 패키지 확인"),
    ("2", "조건 정리", "작동 조건을 절차와 불변식으로 정리"),
    ("3", "AI 실행", "AI + MCP만으로 새 룸 생성 · 빌드 · 검증"),
]
for i, (num, title, desc) in enumerate(cols):
    x = 0.9 + i * 4.05
    rect(sl, x, 2.0, 3.55, 3.15, WHITE, WHITE)
    text(sl, num, x + 0.2, 2.25, 0.7, 0.7, 34, BLUE, True)
    text(sl, title, x + 0.2, 3.05, 3.1, 0.4, 22, BLUE, True)
    text(sl, desc, x + 0.2, 3.72, 3.1, 0.85, 17, BLUE)
text(sl, "환경: Claude Code · Unity MCP 27826 · XRCollabDemo · Unity 6000.3.11f1", 0.9, 6.05, 11.7, 0.35, 17, BLUE)

# 7 C1-C4
sl = add_slide(WHITE)
header(sl, 6, "작동 불변식 C1~C4")
table(sl, [
    ["조건", "내용"],
    ["C1", "서버와 클라이언트의 네트워크 프리팹 컬렉션 일치"],
    ["C2", "플레이어 스폰은 지정 스포너 프리팹과 플랫폼별 아바타 카탈로그 사용"],
    ["C3", "룸 진입/이탈 시 로비 씬과 룸 씬이 정확히 전환"],
    ["C4", "서버는 Master+Room, 에디터는 Client+Room 흐름 기준으로 검증"],
], 0.9, 1.75, 11.55, 4.55, [0.16, 0.84], font_size=13)
text(sl, "의미: 암묵지를 AI가 판정 가능한 조건으로 외부화", 0.9, 6.62, 11.6, 0.35, 20, BLUE, True)

# 8 Result server
sl = add_slide(BLUE)
text(sl, "07", 0.65, 0.55, 0.7, 0.38, 16, WHITE, True)
text(sl, "검증 결과 — 서버와 입장", 1.45, 0.46, 10.8, 0.55, 28, WHITE, True)
line(sl, 0.65, 1.22, 12.05, 0.018, WHITE)
bullets(sl, [
    "`Online Scene: BasicRoom`",
    "`Client 0 is successfully validated`",
    "마스터 서버 192.168.50.49:5000 listening",
    "룸 서버 :7777에서 방 등록",
    "Room-1A5C-ED20-3935 / ID 0 / 최대 10명 / Public",
], 1.1, 1.9, 11.2, 3.3, 22, WHITE)
text(sl, "의미: 룸 생성, 서버 등록, 클라이언트 입장 흐름이 실제 동작", 1.1, 6.35, 11.2, 0.45, 20, WHITE, True)

# 9 Result avatar
sl = add_slide(WHITE)
header(sl, 8, "검증 결과 — 씬 전환과 아바타")
two_columns(
    sl,
    "씬 전환",
    ["Client 씬 언로드", "BasicRoom 로드", "로비 UI C-MasterCanvas 제거", "별도 수동 처리 없이 룸 입장 흐름 유지"],
    "아바타 스폰",
    ["NetworkObjects=5", "Desktop(Clone)이 scene=BasicRoom에 위치", "mixamorig:Hips와 플랫폼 루트도 BasicRoom", "아바타 카메라 활성, 게임뷰에 룸 표시"],
)

# 10 Proof map
sl = add_slide(WHITE)
header(sl, 9, "증명 매핑 — 기준과 증거")
table(sl, [
    ["AI-ready 판단 기준", "실제 증거", "의미"],
    ["AI가 룸을 만들 수 있는가", "MCP로 BasicRoom 생성/빌드", "절차가 실행 가능"],
    ["서버에 등록되는가", "Online Scene, Room registered", "XumLobby 흐름 동작"],
    ["입장이 검증되는가", "Client 0 validated", "인증/입장 연결"],
    ["씬 전환되는가", "Client 언로드, BasicRoom 로드", "로비→룸 전환 성공"],
    ["아바타가 스폰되는가", "Desktop(Clone) in BasicRoom", "XumNet/Motion 연결"],
    ["AI가 판정 가능한가", "C1~C4로 확인", "검증 가능한 절차"],
], 0.45, 1.55, 12.45, 5.75, [0.35, 0.34, 0.31], font_size=10)

# 11 Conclusion
sl = add_slide(CREAM)
header(sl, 10, "결론")
big_quote(sl, "OXR-SDK는 AI-ready한 방식으로 제공될 수 있다")
bullets(sl, [
    "OXR-SDK 계열 프레임워크는 조립 가능한 레이어를 갖고 있음",
    "작동 조건을 문서 · 프롬프트 · 검증 기준으로 외부화함",
    "AI가 MCP로 룸 생성 · 빌드 · 서버 등록 · 아바타 스폰까지 검증함",
    "XRCollabDemo는 OXR-SDK 검증의 대표 샘플이자 best sample",
], 1.05, 4.75, 11.3, 1.75, 17)

# 12 Boundary
sl = add_slide(WHITE)
header(sl, 11, "경계와 해석")
two_columns(
    sl,
    "이번 발표가 말하는 것",
    ["룸 생성 · 입장 · 아바타 스폰 범위 검증", "기존 프레임워크를 AX 관점에서 절차화", "XRCollabDemo를 재생성 가능한 레퍼런스로 전환"],
    "이번 발표가 말하지 않는 것",
    ["기존 문서가 처음부터 완벽했다", "문장 하나로 임의의 룸을 완전 합성한다", "콘텐츠 기능 생성까지 이미 끝났다"],
)

# 13 Roadmap
sl = add_slide(WHITE)
header(sl, 12, "향후 계획")
table(sl, [
    ["Phase", "내용", "상태"],
    ["0", "룸 베이스 검증: 네트워크, 아바타, UI 전환", "완료"],
    ["1", "룸 조립 계약: C1~C4, 스포너 규칙, 씬 토폴로지", "완료"],
    ["2", "AI + MCP 기반 룸 생성 · 빌드 · 검증", "완료"],
    ["3", "콘텐츠 확장 파일럿: IToggleableContent, Ruler", "완료"],
    ["4", "런치패드 UI: 레지스트리 → 아이콘 그리드 → 토글", "예정"],
    ["5", "/RoomContent-<feature> 스킬화", "예정"],
    ["6", "/promptscene + 자동 검증 하네스", "예정"],
], 0.65, 1.55, 12.05, 5.65, [0.14, 0.68, 0.18], font_size=10)

# 14 Closing
sl = add_slide(BLUE)
text(sl, "AI-ready 전", 1.0, 1.25, 5.1, 0.45, 24, WHITE, True)
text(sl, '"XRCollabDemo 샘플은 작동한다"', 1.0, 1.95, 11.2, 0.8, 34, WHITE, True)
line(sl, 1.0, 3.18, 11.3, 0.035, WHITE)
text(sl, "AI-ready 후", 1.0, 3.75, 5.1, 0.45, 24, WHITE, True)
text(sl, '"OXR-SDK 위에서 AI가 협업 룸을 재생성하고 검증할 수 있다"', 1.0, 4.45, 11.3, 1.0, 32, WHITE, True)
text(sl, "XRCollabDemo는 OXR-SDK의 재생성 가능한 best sample", 1.0, 6.45, 11.2, 0.45, 20, WHITE)


out = r"C:\J_0\docs\OXR-SDK-AI-ready-claim-proof.pptx"
prs.save(out)
print(out)
print(f"{len(prs.slides)} slides")
