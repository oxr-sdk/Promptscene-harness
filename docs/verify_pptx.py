from pptx import Presentation
prs = Presentation(r"c:\J_0\docs\OXR-SDK-AI-ready.pptx")
print("Slides:", len(prs.slides))
for i, sl in enumerate(prs.slides, 1):
    texts = [sh.text_frame.text[:55] for sh in sl.shapes
             if sh.has_text_frame and sh.text_frame.text.strip()]
    first = texts[0] if texts else "(empty)"
    print(f"  {i:2d}: {first}")
