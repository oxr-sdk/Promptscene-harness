# -*- coding: utf-8 -*-
"""
OXR-SDK AI-ready  Presentation Generator  v3
Design: Pretendard + Royal Blue #1E40A8 + Cream #F2F0E8 + Flat Modern
12 slides : Cover | Sec01 | x2 body | Sec02 | x2 body |
            Sec03 | x2 body | Sec04+result | Sec05+roadmap

Evidence: docs/*.md  (모든 수치·로그는 실측값)
Font NOTE: Pretendard (OFL, free)  -- 발표 PC 사전 설치 필요
  https://github.com/orioncactus/pretendard/releases
  python-pptx does NOT embed fonts natively;
  install Pretendard on the presenting machine.
"""
import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ── Design Tokens ────────────────────────────────────────────────────────────
BLUE     = RGBColor(0x1E, 0x40, 0xA8)   # Royal Blue  primary
CREAM    = RGBColor(0xF2, 0xF0, 0xE8)   # Cream/Ivory  section bg
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_D   = RGBColor(0x33, 0x36, 0x3D)   # body dark
TEXT_S   = RGBColor(0x8A, 0x8D, 0x94)   # secondary
BLUE_S   = RGBColor(0xC2, 0xCF, 0xEC)   # blue-bg secondary copy
GREEN    = RGBColor(0x16, 0x8A, 0x4E)   # verified
AMBER    = RGBColor(0xD9, 0x77, 0x06)   # warning / roadmap
RED_TXT  = RGBColor(0xB9, 0x1C, 0x1C)
RED_BG   = RGBColor(0xFD, 0xF0, 0xF0)   # fail case bg
GRN_BG   = RGBColor(0xF0, 0xFD, 0xF4)   # success bg
BLUE_BG  = RGBColor(0xEF, 0xF4, 0xFF)   # light blue card
RULE     = RGBColor(0xE4, 0xE7, 0xED)   # hairline

FONT     = 'Pretendard'
W, H     = Inches(13.33), Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]

_log = []

# ── Primitives ───────────────────────────────────────────────────────────────
def new_slide(label=''):
    sl = prs.slides.add_slide(BLANK)
    _log.append((len(prs.slides), label))
    return sl

def bg(sl, color):
    b = sl.background.fill; b.solid(); b.fore_color.rgb = color

def rect(sl, x, y, w, h, fill=WHITE, line=None):
    s = sl.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line: s.line.color.rgb = line; s.line.width = Pt(0.75)
    else:    s.line.fill.background()
    return s

def hrule(sl, y, x1=0.6, x2=12.73, c=RULE):
    s = sl.shapes.add_shape(1, Inches(x1), Inches(y),
                            Inches(x2-x1), Inches(0.015))
    s.fill.solid(); s.fill.fore_color.rgb = c; s.line.fill.background()

def vrule(sl, x, y1, y2, c=BLUE):
    s = sl.shapes.add_shape(1, Inches(x), Inches(y1),
                            Inches(0.015), Inches(y2-y1))
    s.fill.solid(); s.fill.fore_color.rgb = c; s.line.fill.background()

def circle(sl, cx, cy, r, fill=WHITE):
    s = sl.shapes.add_shape(9, Inches(cx-r), Inches(cy-r),
                            Inches(r*2), Inches(r*2))
    s.fill.solid(); s.fill.fore_color.rgb = fill; s.line.fill.background()

def tb(sl, txt, x, y, w, h, size=14, bold=False, color=TEXT_D,
       align=PP_ALIGN.LEFT, italic=False, wrap=True):
    box = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.word_wrap = wrap
    tf = box.text_frame; tf.word_wrap = wrap
    p  = tf.paragraphs[0]; p.alignment = align
    r  = p.add_run(); r.text = txt
    _font(r, size, bold, italic, color)
    return box

def _font(r, size, bold, italic, color):
    r.font.name  = FONT
    r.font.size  = Pt(size)
    r.font.bold  = bold
    r.font.italic = italic
    r.font.color.rgb = color

def _lspc(p, pct):
    """Set paragraph line spacing to pct (e.g. 170 = 170%)."""
    pPr = p._p.get_or_add_pPr()
    lSpc = etree.SubElement(pPr, qn('a:lSpc'))
    etree.SubElement(lSpc, qn('a:spcPct')).set('val', str(pct * 1000))

def aline(tf, txt, size=13, bold=False, color=TEXT_D,
          align=PP_ALIGN.LEFT, italic=False, sp=0, lh=0):
    """Append a paragraph to text-frame tf."""
    p = tf.add_paragraph(); p.alignment = align
    if sp:  p.space_before = Pt(sp)
    if lh:  _lspc(p, lh)
    r = p.add_run(); r.text = txt
    _font(r, size, bold, italic, color)

# ── Layout templates ─────────────────────────────────────────────────────────
def body_slide(label, sec, title):
    """White body: thin blue top bar + >> accent + section + title."""
    sl = new_slide(label)
    bg(sl, WHITE)
    rect(sl, 0, 0, 13.33, 0.06, fill=BLUE)          # top bar
    rect(sl, 0.55, 0.6, 0.045, 0.5, fill=BLUE)       # >> accent
    tb(sl, sec,   0.68, 0.6,  5.0, 0.28, size=10, color=BLUE)
    tb(sl, title, 0.68, 0.82, 12.0, 0.55, size=22, bold=True, color=BLUE)
    hrule(sl, 7.28, 0, 13.33, RULE)
    return sl

def sec_div(num, title, subtitle='', right_content_fn=None):
    """Cream section divider.
    right_content_fn(sl): callback draws the entire right panel.
      When provided, the title/subtitle are NOT added by sec_div —
      the callback is responsible for adding them.
    """
    sl = new_slide(f'섹션 {num}: {title}')
    bg(sl, CREAM)
    tb(sl, num, 0.65, 1.45, 3.2, 3.5, size=108, bold=True, color=BLUE)
    if right_content_fn:
        vrule(sl, 3.85, 1.45, 7.05, BLUE)   # full-height rule
        right_content_fn(sl)                  # callback owns the right zone
    else:
        vrule(sl, 3.85, 1.95, 5.55, BLUE)
        tb(sl, title, 4.1, 2.38, 8.6, 1.1, size=38, bold=True, color=BLUE)
        if subtitle:
            tb(sl, subtitle, 4.1, 3.68, 8.6, 0.55, size=15, color=TEXT_S)
    return sl

# ═════════════════════════════════════════════════════════════════════════════
# S01  Cover
# ═════════════════════════════════════════════════════════════════════════════
sl = new_slide('표지')
bg(sl, BLUE)
circle(sl, 1.2, 1.25, 0.2, WHITE)        # small white circle accent
tb(sl, 'is OXR-SDK', 1.9, 1.75, 9.5, 1.45, size=62, bold=True, color=WHITE)
tb(sl, 'AI ready?',  1.9, 3.12, 9.5, 1.45, size=62, bold=True, color=WHITE)
hrule(sl, 4.88, 1.9, 12.2, WHITE)
tb(sl, ('프레임워크의 AI 준비도 검증 '
        '— XRCollabDemo PromptScene 자율 빌드·검증 결과 보고'),
   1.9, 5.08, 10.2, 0.52, size=14, color=BLUE_S)
tb(sl, 'OXR Research  ·  KISTI  ·  2026-06-16',
   7.5, 7.0, 5.6, 0.34, size=10, color=BLUE_S, align=PP_ALIGN.RIGHT)

# ═════════════════════════════════════════════════════════════════════════════
# S02  Section 01 — 추진배경
# ═════════════════════════════════════════════════════════════════════════════
sec_div('01', '추진배경', 'AI-ready 정의 및 패러다임 전환')

# ═════════════════════════════════════════════════════════════════════════════
# S03  DX→AX + AI-ready 정의
# ═════════════════════════════════════════════════════════════════════════════
sl = body_slide('DX→AX 전환 및 AI-ready 정의',
                '01  추진배경',
                'DX → AX : 새로운 병목,  새로운 기준')

# Left card
rect(sl, 0.55, 1.5, 5.85, 5.65, fill=BLUE_BG)
tb(sl, 'DX 시대  →  AX 시대', 0.75, 1.62, 5.45, 0.4,
   size=12, bold=True, color=BLUE)
hrule(sl, 2.1, 0.75, 6.27, BLUE)

left_box = sl.shapes.add_textbox(Inches(0.75), Inches(2.2), Inches(5.45), Inches(4.7))
left_box.word_wrap = True
tf_l = left_box.text_frame; tf_l.word_wrap = True

aline(tf_l, 'DX 시대의 병목', 13, bold=True, color=AMBER)
aline(tf_l, '"사람이 이 코드를 짤 수 있나?"', 13, color=TEXT_D, lh=170)
aline(tf_l, '', 7)
aline(tf_l, '▼  패러다임 전환  ▼', 13, bold=True, color=BLUE,
      align=PP_ALIGN.CENTER, sp=4)
aline(tf_l, '', 7)
aline(tf_l, 'AX 시대의 병목', 13, bold=True, color=GREEN, sp=4)
aline(tf_l, '"AI 에이전트가 이 프레임워크를', 13, color=TEXT_D, lh=170)
aline(tf_l, '  성공적으로 쓸 수 있나?"', 13, color=TEXT_D)
aline(tf_l, '', 6)
aline(tf_l, 'AI-ready = 경쟁 필수조건', 14, bold=True, color=BLUE, sp=6)

# Right card (blue)
rect(sl, 6.75, 1.5, 6.05, 5.65, fill=BLUE)
tb(sl, 'AI-ready  정의', 6.95, 1.62, 5.65, 0.38, size=13, bold=True, color=WHITE)
tb(sl, '발표 전체의 평가 기준선', 6.95, 2.05, 5.65, 0.35, size=11, color=BLUE_S)
hrule(sl, 2.47, 6.95, 12.62, BLUE_S)

right_box = sl.shapes.add_textbox(
    Inches(6.95), Inches(2.6), Inches(5.65), Inches(4.22))
right_box.word_wrap = True
tf_r = right_box.text_frame; tf_r.word_wrap = True

aline(tf_r, '"SDK의 문서+패키지만으로,', 15, bold=True, color=WHITE, lh=170)
aline(tf_r, '  사람의 암묵지 없이,',     15, bold=True, color=WHITE, lh=170)
aline(tf_r, '  AI 에이전트가',           15, bold=True, color=WHITE, lh=170)
aline(tf_r, '  검증 가능하게 작동하는',  15, bold=True, color=BLUE_S, lh=170)
aline(tf_r, '  산출물을 만들 수 있는 상태."',  15, bold=True, color=WHITE, lh=170)
aline(tf_r, '', 8)
aline(tf_r, "핵심:  '그럴듯함'  ≠  '증명 가능함'",
      12, italic=True, color=BLUE_S, sp=5, lh=160)
aline(tf_r, '→ 이후 모든 주장의 기준선',
      12, italic=True, color=BLUE_S, sp=4, lh=160)

# ═════════════════════════════════════════════════════════════════════════════
# S04  클론 vs 프롬프트 재생성
# ═════════════════════════════════════════════════════════════════════════════
sl = body_slide('클론 vs 프롬프트 재생성',
                '01  추진배경',
                '정적 클론 vs 프롬프트 재생성 — 왜 다른가')

# Table header row
rect(sl, 0.55, 1.5, 12.25, 0.48, fill=BLUE)
tb(sl, '항목',    0.72, 1.58, 2.75, 0.32, size=12, bold=True, color=WHITE)
tb(sl, '정적 클론',     3.55, 1.58, 4.15, 0.32, size=12, bold=True, color=WHITE)
tb(sl, '프롬프트 재생성 (PromptScene)',
   7.82, 1.58, 4.8, 0.32, size=12, bold=True, color=WHITE)

TABLE_ROWS = [
    ('지식 형태',     '동결 스냅샷',                  '문서/계약(C1~C4)에 노출'),
    ('크기',          '전체 수GB 클론',               '핵심 문서 + 패키지'),
    ('조합',          '불가 — 통째로만',              '조합 가능 (콘텐츠 모듈화)'),
    ('SDK 업데이트',  'drift 발생',                   '문서 갱신으로 동기화'),
    ('암묵지',        '은닉 (클론 안에 숨어있음)',     '구조적 검증으로 제거'),
    ('검증 기준',     '동일 파일 여부',               'C1~C4 불변식 충족 여부'),
]
for i, (item, clone_v, prompt_v) in enumerate(TABLE_ROWS):
    yy  = 2.05 + i * 0.65
    bg_c = WHITE if i % 2 == 0 else BLUE_BG
    rect(sl, 0.55, yy, 12.25, 0.62, fill=bg_c)
    tb(sl, item,     0.72, yy+0.12, 2.75, 0.4, size=11, bold=True, color=TEXT_D)
    tb(sl, clone_v,  3.55, yy+0.12, 4.15, 0.4, size=11, color=AMBER)
    tb(sl, prompt_v, 7.82, yy+0.12, 4.8,  0.4, size=11, color=GREEN)

# Callout — spawner clone discovery
rect(sl, 0.55, 6.1, 12.25, 1.05, fill=RGBColor(0xFF, 0xFB, 0xEB), line=AMBER)
tb(sl, '실증: 스포너 클론 발견 — 클론은 작동 이유를 숨긴다',
   0.75, 6.18, 11.8, 0.38, size=13, bold=True, color=AMBER)
tb(sl, ("DocRoom이 '됐던' 것은 Room.unity 스포너를 몰래 클론했기 때문."
        " 순수 문서만으론 안 됐음."
        " → 문서+프리팹(Room-PlayerSpawner.prefab)으로 메꿔 진짜 재현 가능하게 함."),
   0.75, 6.58, 11.8, 0.5, size=11, color=TEXT_D)

# ═════════════════════════════════════════════════════════════════════════════
# S05  Section 02 — 추진내용
# ═════════════════════════════════════════════════════════════════════════════
sec_div('02', '추진내용', '자율 빌드 절차 및 콘텐츠 플러그인 아키텍처')

# ═════════════════════════════════════════════════════════════════════════════
# S06  자율 빌드 절차
# ═════════════════════════════════════════════════════════════════════════════
sl = body_slide('자율 빌드 절차', '02  추진내용',
                '패키지 문서 → MCP 자율 빌드 → 서버 빌드 → 구조적 검증')

STEPS = [
    ('①', '패키지 문서 수집',
     'OXR-SDK README + GitBook 공식 문서\n→ c:\\J_0\\docs\\*.md 하나로 정리'),
    ('②', '자율 빌드  (AI + MCP)',
     'Claude Code + Unity MCP (포트 27826)\n→ 씬·스크립트·설정 자율 적용'),
    ('③', '서버 빌드',
     'XumLobbyServerBuilderWindow reflection 호출\n→ Master.exe + Room.exe 생성'),
    ('④', '구조적 검증',
     'C1~C4 불변식 + 런타임 로그 패턴 체크\n→ 전부 통과 = AI-ready 판정'),
]
for i, (num, title_s, desc) in enumerate(STEPS):
    xx = 0.38 + i * 3.24
    rect(sl, xx, 1.52, 3.1, 5.55, fill=BLUE_BG)
    rect(sl, xx, 1.52, 3.1, 0.55, fill=BLUE)   # top banner
    tb(sl, num,    xx+0.12, 1.57, 0.5,  0.42, size=18, bold=True, color=WHITE)
    tb(sl, title_s, xx+0.6, 1.6,  2.35, 0.42, size=13, bold=True, color=WHITE)
    tb(sl, desc, xx+0.14, 2.22, 2.82, 4.6,  size=12, color=TEXT_D, wrap=True)
    if i < 3:
        tb(sl, '→', xx+3.0, 4.1, 0.32, 0.42, size=18, bold=True,
           color=BLUE, align=PP_ALIGN.CENTER)

tb(sl, ('XRCollabDemo는 증거(evidence)다 '
        '— 주인공은 OXR-SDK + PromptScene 프레임워크'),
   0.55, 7.08, 12.25, 0.34, size=11, italic=True,
   color=TEXT_S, align=PP_ALIGN.CENTER)

# ═════════════════════════════════════════════════════════════════════════════
# S07  콘텐츠 플러그인 아키텍처
# ═════════════════════════════════════════════════════════════════════════════
sl = body_slide('콘텐츠 플러그인 아키텍처', '02  추진내용',
                '씬 계층 표준 + 조합 가능성 증거 — Ruler 파일럿 (Phase 2 완료)')

# Left — scene hierarchy
rect(sl, 0.55, 1.5, 6.05, 5.65, fill=BLUE_BG)
tb(sl, '씬 계층 표준 (Phase 2.5 완료)', 0.75, 1.62, 5.65, 0.38,
   size=12, bold=True, color=BLUE)
hrule(sl, 2.08, 0.75, 6.43, BLUE)

LAYERS = [
    ('SYSTEMS',      'Network · PlayerSpawner · RoomCore',  BLUE,  2.18),
    ('ENVIRONMENT',  'Floor · Walls · Lighting · Camera',    TEXT_S, 3.0),
    ('UI',           'R-MasterCanvas · Launchpad (Phase 3)', TEXT_S, 3.82),
    ('FEATURES',     'Ruler  ←  IToggleableContent  ✓ 실측',  GREEN, 4.64),
    ('_DYNAMIC',     '아바타 Clone · 측정선 (런타임 생성)',     TEXT_S, 5.46),
]
for name, desc, col, yy in LAYERS:
    rect(sl, 0.72, yy, 5.7, 0.72, fill=WHITE, line=col)
    tb(sl, f'===== {name} =====', 0.85, yy+0.05, 5.45, 0.3,
       size=11, bold=True, color=col)
    tb(sl, desc, 0.85, yy+0.35, 5.45, 0.3, size=10, color=TEXT_S)

# Right top — principle
rect(sl, 6.95, 1.5, 5.85, 2.82, fill=GRN_BG, line=GREEN)
tb(sl, '핵심 원칙 (UE5 Modular Game Features 동일)',
   7.12, 1.62, 5.5, 0.38, size=12, bold=True, color=GREEN)
hrule(sl, 2.08, 7.12, 12.63, GREEN)

rp = sl.shapes.add_textbox(Inches(7.12), Inches(2.17), Inches(5.5), Inches(1.9))
rp.word_wrap = True
tf_p = rp.text_frame; tf_p.word_wrap = True
aline(tf_p, '의존 방향: FEATURE → SYSTEMS (단방향)', 12, color=TEXT_D, lh=165)
aline(tf_p, 'SYSTEMS는 특정 FEATURE를 컴파일타임에 모름', 12, color=TEXT_S, lh=165)
aline(tf_p, "판별: '빼도 안 깨지면 FEATURE'", 12, italic=True, color=TEXT_S, lh=165)
aline(tf_p, '→ 런타임 on/off + 프로젝트 간 이식 가능', 12, bold=True, color=GREEN, sp=4)

# Right bottom — Ruler pilot
rect(sl, 6.95, 4.62, 5.85, 2.53, fill=BLUE_BG, line=BLUE)
tb(sl, 'Ruler 콘텐츠 파일럿 (Phase 2 완료 — 실측)',
   7.12, 4.74, 5.5, 0.38, size=12, bold=True, color=BLUE)
hrule(sl, 5.2, 7.12, 12.63, BLUE)

rr = sl.shapes.add_textbox(Inches(7.12), Inches(5.28), Inches(5.5), Inches(1.7))
rr.word_wrap = True
tf_rr = rr.text_frame; tf_rr.word_wrap = True
aline(tf_rr, 'RulerContent : IToggleableContent', 12, bold=True, color=TEXT_D, lh=165)
aline(tf_rr, 'RoomCore.Instance에 자기등록 → SetEnabled(true)', 11, color=TEXT_S, lh=165)
aline(tf_rr, '두 점 측정선 + 거리 라벨:  4.00 m  (실측)', 12, bold=True,
      color=GREEN, sp=3)
aline(tf_rr, 'DeepChairProject 앱레이어 의존 = 0 → 클린 재구현',
      11, italic=True, color=TEXT_S, sp=3)

# ═════════════════════════════════════════════════════════════════════════════
# S08  Section 03 — 검증결과
# ═════════════════════════════════════════════════════════════════════════════
sec_div('03', '검증결과', '작동 불변식 C1~C4 실측 발굴 및 런타임 증거')

# ═════════════════════════════════════════════════════════════════════════════
# S09  C1~C4 불변식 + 스포너 갭
# ═════════════════════════════════════════════════════════════════════════════
sl = body_slide('작동 불변식 C1~C4', '03  검증결과',
                '이 4개 중 하나라도 위반하면 룸이 안 돔 — 전부 실측으로 발굴')

CONTRACTS = [
    ('C1', '프리팹 컬렉션 일치',
     '서버·클라 NetworkManager._spawnablePrefabs = DefaultPrefabObjects (동일)',
     '위반 시:  입장은 되나 아바타 안 보임'),
    ('C2', '플레이어 스포너',
     'XumPlayerSpawner + Catalog(Desktop/UnityXR) — 반드시 프리팹 인스턴스화\n'
     'Room-PlayerSpawner.prefab  (Assets/PromptScene/Prefabs/)',
     "'Failed to confirm the access'  →  즉시 퇴장"),
    ('C3', '씬 전환',
     "R-RoomServer.DefaultScene: _onlineScene=<룸경로> / _offlineScene=Client.unity",
     '위반 시:  로비 UI가 룸 위를 덮음'),
    ('C4', '실행 토폴로지',
     '서버 = Master.exe + Room.exe (둘 다 필수)  /  에디터 = Client + Room 동시 로드',
     '위반 시:  네트워크 씬 전환 미발동'),
]
for i, (cid, name, spec, fail) in enumerate(CONTRACTS):
    yy = 1.52 + i * 1.38
    rect(sl, 0.55, yy, 12.25, 1.28, fill=WHITE, line=RULE)
    # C-badge (blue left strip)
    rect(sl, 0.55, yy, 0.82, 1.28, fill=BLUE)
    tb(sl, cid, 0.56, yy+0.38, 0.8, 0.52, size=20, bold=True,
       color=WHITE, align=PP_ALIGN.CENTER)
    tb(sl, name, 1.52, yy+0.07, 3.95, 0.42, size=13, bold=True, color=TEXT_D)
    tb(sl, spec, 1.52, yy+0.52, 5.42, 0.72, size=10, color=TEXT_S, wrap=True)
    # fail case
    rect(sl, 7.12, yy+0.14, 5.55, 1.0, fill=RED_BG)
    tb(sl, fail, 7.25, yy+0.32, 5.3, 0.62, size=11, color=RED_TXT, wrap=True)
    tb(sl, '✓ 실측', 11.82, yy+0.07, 0.88, 0.35, size=11, bold=True, color=GREEN)

# ═════════════════════════════════════════════════════════════════════════════
# S10  런타임 실측 (BasicRoom)
# ═════════════════════════════════════════════════════════════════════════════
sl = body_slide('런타임 실측 — BasicRoom 검증', '03  검증결과',
                '근거: PromptScene-progress-report.md §3-1  /  build-xumlobby-server.md 부록')

# 3-column log cards
LOG_COLS = [
    ('서버 측 로그', BLUE, [
        ('Online Scene: BasicRoom',               '씬 등록 ✓'),
        ('Room registered successfully.',          'Room-1A5C-ED20-3935'),
        ('Client 0 is successfully validated',    '클라이언트 검증 ✓'),
        ('192.168.50.49:5000  /  :7777 (UDP)',    '실측 IP·포트'),
        ('빌드 시간: ~48초 (캐시 상태)',             'Master+Room 합산'),
    ]),
    ('에디터 클라이언트', RGBColor(0x11, 0x72, 0xB8), [
        ('Client 씬 언로드됨',                    'C3 로비 자동 소멸 ✓'),
        ('C-MasterCanvas 소멸',                   'UI 자동 전환 ✓'),
        ('NetworkObjects = 5',                    '실측 오브젝트 수'),
        ('Desktop(Clone)  owner=True',            '아바타 스폰 ✓'),
        ('scene = BasicRoom',                     'C1~C3 토폴로지 ✓'),
    ]),
    ('기능 검증', GREEN, [
        ('아바타 카메라 활성',                     '게임뷰 정상'),
        ('게임뷰에 룸 표시 (로비 없음)',             'C3 완전 충족'),
        ('WASD 이동 확인',                         '로코모션 동작'),
        ('Ruler 측정선:  4.00 m',                 '콘텐츠 플러그인 ✓'),
        ('SetEnabled(true) 예외 없음',             'IToggleableContent ✓'),
    ]),
]
for ci, (title_c, col, items) in enumerate(LOG_COLS):
    xx = 0.38 + ci * 4.35
    rect(sl, xx, 1.5, 4.15, 5.9, fill=WHITE, line=col)
    rect(sl, xx, 1.5, 4.15, 0.48, fill=col)
    tb(sl, title_c, xx+0.16, 1.57, 3.83, 0.35, size=13, bold=True, color=WHITE)
    for ri, (log_txt, note) in enumerate(items):
        yy = 2.05 + ri * 1.06
        rect(sl, xx+0.16, yy, 3.83, 0.95, fill=BLUE_BG)
        tb(sl, log_txt, xx+0.26, yy+0.06, 3.62, 0.42,
           size=10, bold=True, color=TEXT_D, wrap=True)
        tb(sl, note,    xx+0.26, yy+0.58, 3.62, 0.3,
           size=9, italic=True, color=col)

# ═════════════════════════════════════════════════════════════════════════════
# S11  Section 04 — 결론  (divider + verdict combined on blue)
# ═════════════════════════════════════════════════════════════════════════════
def _result_content(sl):
    """Right panel for slide 11 (결론).
    Owns the full right zone x>4.1.  sec_div does NOT add a title here.
    Layout:
      y=1.52  Section title "결론"
      y=2.12  hrule
      y=2.22  Verdict box  (h=1.78)
      y=4.18  Two columns: proven (left) / roadmap (right)  (h=2.97)
    """
    # Section title
    tb(sl, '결론', 4.1, 1.52, 8.95, 0.55, size=28, bold=True, color=BLUE)
    hrule(sl, 2.14, 4.1, 13.1, RULE)

    # Verdict box
    rect(sl, 4.1, 2.22, 9.0, 1.78, fill=GRN_BG, line=GREEN)
    rect(sl, 4.1, 2.22, 0.08, 1.78, fill=GREEN)
    tb(sl, '✓  OXR-SDK는 AI-ready 정의를 충족한다',
       4.35, 2.38, 8.5, 0.55, size=17, bold=True, color=BLUE)
    tb(sl, ('AI 에이전트가 SDK 문서만 보고,'
            ' 검증 가능하게 작동하는 룸을 자율 빌드·검증했다.'
            ' (BasicRoom 실측 완료)'),
       4.35, 2.96, 8.5, 0.85, size=11, color=TEXT_D, wrap=True)

    # Proven list (left column)
    rect(sl, 4.1, 4.18, 4.35, 2.97, fill=GRN_BG, line=GREEN)
    tb(sl, '증명선 (Phase 0~2.5)', 4.28, 4.3, 4.0, 0.38,
       size=12, bold=True, color=GREEN)
    PROVED = [
        '✓ BasicRoom 자율 빌드·검증 완료',
        '✓ C1~C4 불변식 전부 실측 발굴',
        '✓ 스포너 클론 갭 → 프리팹 해소',
        '✓ Ruler 4.00 m 콘텐츠 플러그인',
        '✓ 서버 빌드 자동화 (reflection)',
    ]
    for i, item in enumerate(PROVED):
        tb(sl, item, 4.28, 4.8+i*0.48, 4.0, 0.42, size=11, color=TEXT_D)

    # Roadmap list (right column)
    rect(sl, 8.62, 4.18, 4.45, 2.97, fill=BLUE_BG, line=BLUE)
    tb(sl, '로드맵선 (Phase 3~5)', 8.8, 4.3, 4.08, 0.38,
       size=12, bold=True, color=BLUE)
    ROAD = [
        ('Phase 3', '런치패드 UI'),
        ('Phase 4', '/RoomContent-* 스킬화'),
        ('Phase 5', '자연어 → 룸 합성'),
        ('', ''),
        ('⚠', '자율 빌드 ≠ 자유 합성'),
    ]
    for i, (ph, desc) in enumerate(ROAD):
        col_t = AMBER if '⚠' in ph else (BLUE if ph else TEXT_S)
        txt = f'{ph}  {desc}' if ph else ''
        tb(sl, txt, 8.8, 4.8+i*0.48, 4.08, 0.42, size=11, color=col_t)

sec_div('04', '결론', right_content_fn=_result_content)

# ═════════════════════════════════════════════════════════════════════════════
# S12  Section 05 — 향후계획  (divider + roadmap combined on cream)
# ═════════════════════════════════════════════════════════════════════════════
def _roadmap_content(sl):
    """Right panel for slide 12 (향후계획).
    Owns the full right zone x>4.1.  sec_div does NOT add a title here.
    Layout:
      y=1.52  Section title "향후계획"
      y=2.12  hrule
      y=2.22  Done-phase chips  (h=0.48)
      y=2.88  Three Phase cards  (h=4.07)
      y=7.02  Bottom note
    """
    # Section title
    tb(sl, '향후계획', 4.1, 1.52, 8.95, 0.55, size=28, bold=True, color=BLUE)
    hrule(sl, 2.14, 4.1, 13.1, RULE)

    # Done chips
    DONE = [('Phase 0', '룸 베이스'), ('Phase 1', '계약 규격'),
            ('Phase 2', 'Ruler 파일럿'), ('2.5', '씬 표준화')]
    rect(sl, 4.1, 2.22, 9.0, 0.48, fill=GRN_BG)
    tb(sl, '완료 :', 4.28, 2.3, 0.78, 0.3, size=11, bold=True, color=GREEN)
    for i, (ph, nm) in enumerate(DONE):
        tb(sl, f'✓ {ph}  {nm}', 5.12+i*1.88, 2.3, 1.8, 0.3,
           size=10, color=GREEN)

    # Next 3 Phase cards
    # card width = (13.1-4.1) / 3 - 0.1 gap = 2.9"
    PHASES = [
        ('3', '런치패드 UI', 'Runtime',
         ['registry → 아이콘 그리드 자동 생성',
          '스마트폰식 SetEnabled 토글',
          '룸 안에서 기능을 on/off'], BLUE, '개발 중'),
        ('4', '스킬화', 'Authoring',
         ['/RoomContent-<feature> 명령어',
          'LLM 신규기능 생성 템플릿',
          'AI가 새 콘텐츠 모듈을 생성'], RGBColor(0x11, 0x72, 0xB8), '계획'),
        ('5', '자연어→룸 합성', 'Authoring+Runtime',
         ['/promptscene : 자연어 → 룸 조립',
          'C1~C4 + 콘텐츠 자동 검증 하네스',
          'DeepChairProject 기능 자연어 이식'], AMBER, '로드맵'),
    ]
    CW = 2.88   # card width
    for i, (num, title_p, layer, items, col, stat) in enumerate(PHASES):
        xx = 4.1 + i * (CW + 0.12)
        rect(sl, xx, 2.88, CW, 4.07, fill=WHITE, line=col)
        rect(sl, xx, 2.88, CW, 0.5, fill=col)         # banner
        tb(sl, f'Phase {num}', xx+0.12, 2.94, 1.6, 0.38,
           size=12, bold=True, color=WHITE)
        tb(sl, stat, xx+1.75, 2.94, 1.0, 0.38,
           size=10, color=WHITE, align=PP_ALIGN.RIGHT)
        tb(sl, title_p, xx+0.12, 3.5, CW-0.2, 0.48,
           size=14, bold=True, color=TEXT_D)
        tb(sl, f'Layer: {layer}', xx+0.12, 4.02, CW-0.2, 0.3,
           size=10, italic=True, color=col)
        hrule(sl, 4.38, xx+0.12, xx+CW-0.08, col)
        for ri, item in enumerate(items):
            tb(sl, f'• {item}', xx+0.16, 4.52+ri*0.75, CW-0.22, 0.68,
               size=10, color=TEXT_D, wrap=True)

    tb(sl, ('핵심 구분  —  '
            'Phase 3~4 : 자율 빌드 심화  |  Phase 5 : 자유 합성 (별도 단계)'),
       4.1, 7.02, 9.0, 0.34, size=11, bold=True,
       color=AMBER, align=PP_ALIGN.CENTER)

sec_div('05', '향후계획', right_content_fn=_roadmap_content)

# ── Save ─────────────────────────────────────────────────────────────────────
OUT = r'c:\J_0\docs\OXR-SDK-AI-ready.pptx'
prs.save(OUT)
print(f'Saved  {OUT}')
print(f'Total  {len(prs.slides)} slides')
for no, lbl in _log:
    print(f'  {no:2d}.  {lbl}')
